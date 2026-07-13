#!/usr/bin/env python3
"""Build <deck>/<name>.pptx from the deck's slides/*.svg as NATIVE, editable shapes.

Translation (everything below is real, click-to-edit PowerPoint, not a picture):
  <rect>        -> rectangle / rounded-rectangle (fill, line, corner radius)
  <circle>      -> oval
  <line>        -> straight connector
  <text>/<tspan>-> text box (font, size, weight, color, alignment; inherited attrs
                   from ancestor <g> resolved; per-tspan color = separate run)
  gradient bg   -> native gradient-filled rectangle

Decorative <path> (icons, arrows, checks, shield) and faint <g opacity=...>
decorations are composited into ONE transparent full-canvas PNG overlaid on top
(needs an svg renderer; auto-detected). They stay pixel-accurate but flat.

Requires:  python-pptx, lxml          (pip install python-pptx lxml)
Optional (for the icon overlay): Chrome/Edge | rsvg-convert | inkscape | cairosvg | LibreOffice

Usage:   python build_pptx.py <deck-dir> [output.pptx]
Windows: run with  PYTHONUTF8=1 .
"""
import sys, os, re, glob, copy, shutil, subprocess, tempfile
from lxml import etree

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

SVG = "http://www.w3.org/2000/svg"
EMUPX = 9525          # EMU per SVG px (96 dpi);  1280px -> 12192000 EMU (13.333in)
PT = 0.75             # px -> pt
BASELINE = 0.86       # text baseline offset above box top, in em (calibrated)
NAMED = {"white": "FFFFFF", "black": "000000", "red": "FF0000"}
# native-translated elements; polyline/polygon are NOT here on purpose —
# they fall through to the icon overlay instead of silently vanishing
SHAPELIKE = ("rect", "circle", "ellipse", "line", "text")


def E(v):
    return Emu(int(round(float(v) * EMUPX)))


def ln(el):
    t = el.tag
    if not isinstance(t, str):   # comment / processing-instruction node
        return ""
    return etree.QName(t).localname


def color(c):
    """SVG paint -> 'RRGGBB' | None(no paint) | 'GRAD'."""
    if c is None:
        return "MISSING"
    c = c.strip()
    if c == "none":
        return None
    if c.startswith("url("):
        return "GRAD"
    if c.startswith("#"):
        h = c[1:]
        if len(h) == 3:
            h = "".join(ch * 2 for ch in h)
        return h.upper()
    return NAMED.get(c.lower(), "MISSING")


def inh(el, attr, default=None):
    """Resolve a presentation attribute, walking up ancestors."""
    cur = el
    while cur is not None and etree.iselement(cur):
        v = cur.get(attr)
        if v is not None:
            return v
        cur = cur.getparent()
    return default


def in_faint_group(el):
    cur = el.getparent()
    while cur is not None and etree.iselement(cur):
        if ln(cur) == "g" and cur.get("opacity") is not None:
            return True
        cur = cur.getparent()
    return False


def grad_stops(root, gid):
    for lg in root.iter("{%s}linearGradient" % SVG):
        if lg.get("id") == gid:
            return [(float(s.get("offset", 0)), color(s.get("stop-color")))
                    for s in lg if ln(s) == "stop"]
    return None


# ----------------------------- styling helpers -----------------------------
def _fill(sp, paint, root):
    if paint == "GRAD":
        m = None
        # find the url(#id) on the element is handled by caller passing stops
        return
    if paint is None or paint == "MISSING":
        sp.fill.background()
        return
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(paint)


def _line(sp, stroke, width):
    sc = color(stroke) if stroke is not None else None
    if sc and sc not in ("GRAD", "MISSING"):
        sp.line.color.rgb = RGBColor.from_string(sc)
        sp.line.width = E(float(width or 1))
    else:
        sp.line.fill.background()


def _noshadow(sp):
    try:
        sp.shadow.inherit = False
    except Exception:
        pass


def _set_ea(run, name):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)


def _set_spc(run, px):
    run._r.get_or_add_rPr().set("spc", str(int(float(px) * 75)))


def est_width(runs, fs, spc):
    n = 0
    tot = 0.0
    for txt, _, _ in runs:
        for ch in txt:
            n += 1
            tot += fs * (0.55 if ord(ch) < 0x2E80 else 1.0)
    return tot + (spc or 0) * max(0, n - 1)


# ------------------------------- elements ----------------------------------
def add_rect(shapes, el, root):
    x = float(el.get("x", 0)); y = float(el.get("y", 0))
    w = float(el.get("width")); h = float(el.get("height"))
    rx = el.get("rx")
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rx else MSO_SHAPE.RECTANGLE
    sp = shapes.add_shape(shp, E(x), E(y), E(w), E(h))
    if rx:
        try:
            sp.adjustments[0] = max(0.0, min(0.5, float(rx) / min(w, h)))
        except Exception:
            pass
    fill = inh(el, "fill")
    fc = color(fill)
    if fc == "GRAD":
        gid = re.search(r"url\(#([^)]+)\)", fill).group(1)
        stops = grad_stops(root, gid) or [(0, "0B2A4A"), (1, "11406F")]
        sp.fill.gradient()
        try:
            gss = sp.fill.gradient_stops
            for i, (off, col) in enumerate(stops[:len(gss)]):
                gss[i].position = off
                if col and col not in ("GRAD", "MISSING"):
                    gss[i].color.rgb = RGBColor.from_string(col)
            sp.fill.gradient_angle = 45.0
        except Exception:
            pass
    elif fc == "MISSING":
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string("000000")
    else:
        _fill(sp, fc, root)
    _line(sp, inh(el, "stroke"), inh(el, "stroke-width"))
    _noshadow(sp)


def add_circle(shapes, el, root):
    cx = float(el.get("cx")); cy = float(el.get("cy")); r = float(el.get("r"))
    sp = shapes.add_shape(MSO_SHAPE.OVAL, E(cx - r), E(cy - r), E(2 * r), E(2 * r))
    fc = color(inh(el, "fill"))
    _fill(sp, ("000000" if fc == "MISSING" else fc), root)
    _line(sp, inh(el, "stroke"), inh(el, "stroke-width"))
    _noshadow(sp)


def add_ellipse(shapes, el, root):
    cx = float(el.get("cx", 0)); cy = float(el.get("cy", 0))
    rx = float(el.get("rx", 0)); ry = float(el.get("ry", 0))
    sp = shapes.add_shape(MSO_SHAPE.OVAL, E(cx - rx), E(cy - ry), E(2 * rx), E(2 * ry))
    fc = color(inh(el, "fill"))
    _fill(sp, ("000000" if fc == "MISSING" else fc), root)
    _line(sp, inh(el, "stroke"), inh(el, "stroke-width"))
    _noshadow(sp)


def add_line(shapes, el):
    cn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                              E(el.get("x1")), E(el.get("y1")),
                              E(el.get("x2")), E(el.get("y2")))
    sc = color(inh(el, "stroke"))
    if sc and sc not in ("GRAD", "MISSING"):
        cn.line.color.rgb = RGBColor.from_string(sc)
    cn.line.width = E(float(inh(el, "stroke-width") or 1))
    _noshadow(cn)


def add_text(shapes, el):
    x = float(el.get("x", 0)); y = float(el.get("y", 0))
    fs = float(inh(el, "font-size", 16))
    anchor = inh(el, "text-anchor", "start")
    spc = inh(el, "letter-spacing")
    spc = float(spc) if spc else 0
    base_fill = color(inh(el, "fill"))
    base_bold = inh(el, "font-weight") == "bold"
    if base_fill in (None, "GRAD", "MISSING"):
        base_fill = "000000"

    runs = []
    if el.text and el.text.strip():
        runs.append((el.text, base_fill, base_bold))
    for ch in el:
        if ln(ch) == "tspan":
            tf = color(ch.get("fill")) or base_fill
            if tf in ("GRAD", "MISSING"):
                tf = base_fill
            tb = (ch.get("font-weight") == "bold") or base_bold
            if ch.text:
                runs.append((ch.text, tf, tb))
        if ch.tail and ch.tail.strip():
            runs.append((ch.tail, base_fill, base_bold))
    if not runs:
        return

    w = est_width(runs, fs, spc)
    bw = max(w, 4)
    # box edge lands exactly on the svg x for the given anchor; with word_wrap off
    # the line overflows symmetrically, so a rough width estimate is fine.
    left = {"start": x, "middle": x - bw / 2, "end": x - bw}[anchor]
    top = y - fs * BASELINE
    box = shapes.add_textbox(E(left), E(top), E(bw), E(fs * 1.45))
    fr = box.text_frame
    fr.word_wrap = False
    fr.auto_size = MSO_AUTO_SIZE.NONE
    fr.margin_left = fr.margin_right = fr.margin_top = fr.margin_bottom = 0
    fr.vertical_anchor = MSO_ANCHOR.TOP
    p = fr.paragraphs[0]
    p.alignment = {"start": PP_ALIGN.LEFT, "middle": PP_ALIGN.CENTER, "end": PP_ALIGN.RIGHT}[anchor]
    for txt, col, bd in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(fs * PT)
        r.font.bold = bd
        r.font.name = "Microsoft YaHei"
        r.font.color.rgb = RGBColor.from_string(col)
        _set_ea(r, "Microsoft YaHei")
        if spc:
            _set_spc(r, spc)


# ------------------------------ icon overlay -------------------------------
def overlay_svg(svg_text):
    """Strip everything except <path> and faint <g opacity> decorations."""
    root = etree.fromstring(svg_text.encode("utf-8"))
    for el in list(root.iter()):
        if ln(el) in SHAPELIKE and not in_faint_group(el):
            par = el.getparent()
            if par is not None:
                par.remove(el)
    return etree.tostring(root, encoding="unicode")


def find_chrome():
    import shutil as sh
    p = sh.which("chrome") or sh.which("google-chrome") or sh.which("msedge")
    for c in [p, r"C:\Program Files\Google\Chrome\Application\chrome.exe",
              r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"]:
        if c and os.path.isfile(c):
            return c
    return None


def render_overlays(svgs, outdir):
    """name -> transparent overlay png, or {} if no renderer."""
    chrome = find_chrome()
    out = {}
    if chrome:
        prof = tempfile.mkdtemp(prefix="cr_")
        for s in svgs:
            name = os.path.splitext(os.path.basename(s))[0]
            stripped = overlay_svg(open(s, encoding="utf-8").read())
            html = ('<!doctype html><meta charset="utf-8"><style>'
                    'html,body{margin:0;padding:0;background:transparent}'
                    'svg{display:block;width:1280px;height:720px}</style>' + stripped)
            hp = os.path.join(outdir, name + ".html")
            open(hp, "w", encoding="utf-8").write(html)
            png = os.path.join(outdir, name + ".png")
            subprocess.run([chrome, "--headless=new", "--disable-gpu", "--no-sandbox",
                            "--no-first-run", "--no-default-browser-check", "--hide-scrollbars",
                            "--default-background-color=00000000", "--force-device-scale-factor=2",
                            "--window-size=1280,720", "--user-data-dir=" + prof,
                            "--screenshot=" + png, "file:///" + hp.replace("\\", "/")],
                           check=True, capture_output=True)
            os.remove(hp); out[name] = png
        shutil.rmtree(prof, ignore_errors=True)
        return out
    return out


# --------------------------------- build -----------------------------------
def build(svgs, out):
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]

    tmp = tempfile.mkdtemp(prefix="ov_")
    overlays = render_overlays(svgs, tmp)

    for s in svgs:
        name = os.path.splitext(os.path.basename(s))[0]
        root = etree.fromstring(open(s, encoding="utf-8").read().encode("utf-8"))
        slide = prs.slides.add_slide(blank)
        shapes = slide.shapes
        for el in root.iter():
            t = ln(el)
            if t not in SHAPELIKE or in_faint_group(el):
                continue
            if t == "rect":
                add_rect(shapes, el, root)
            elif t == "circle":
                add_circle(shapes, el, root)
            elif t == "ellipse":
                add_ellipse(shapes, el, root)
            elif t == "line":
                add_line(shapes, el)
            elif t == "text":
                add_text(shapes, el)
        if name in overlays:
            shapes.add_picture(overlays[name], 0, 0, prs.slide_width, prs.slide_height)

    prs.save(out)
    shutil.rmtree(tmp, ignore_errors=True)
    return bool(overlays)


def main():
    if len(sys.argv) < 2:
        sys.exit("Usage: build_pptx.py <deck-dir> [output.pptx]")
    deck = sys.argv[1]
    svgs = sorted(glob.glob(os.path.join(deck, "slides", "*.svg")))
    if not svgs:
        sys.exit("No slides/*.svg found in " + deck)
    out = (sys.argv[2] if len(sys.argv) > 2
           else os.path.join(deck, os.path.basename(os.path.abspath(deck)) + ".pptx"))
    had_overlay = build(svgs, out)
    note = "with icon overlay" if had_overlay else "NO renderer -> icons skipped"
    print("Wrote %s (%d slides, native shapes, %s)" % (out, len(svgs), note))


if __name__ == "__main__":
    main()
